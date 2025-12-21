import time
import json
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from backend.state import ResearchState
from backend.utils import call_openai, load_prompt, get_model_settings
from backend.logger import execution_logger


def create_powerpoint(slides_data, topic, output_path):
    """
    Create a PowerPoint presentation from slide data.

    Args:
        slides_data: List of slide dictionaries with title, content, speaker_notes
        topic: Research topic
        output_path: Path to save the .pptx file
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in slides_data:
        slide_type = slide_data.get('slide_type', 'content')

        if slide_type == 'title' and slide_data['slide_number'] == 1:
            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = slide_data['title']

            # Combine content items for subtitle
            subtitle_text = '\n'.join(slide_data['content'][:2]) if slide_data['content'] else ''
            subtitle.text = subtitle_text

        else:
            # Content slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = slide_data['title']

            # Add content to body
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear()

            for idx, point in enumerate(slide_data['content']):
                if idx == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = point
                p.level = 0
                p.font.size = Pt(18)
                p.space_after = Pt(12)

        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = slide_data.get('speaker_notes', '')

    # Save presentation
    prs.save(output_path)
    print(f"✓ PowerPoint saved: {output_path}")


def slide_generation_node(state: ResearchState) -> ResearchState:
    """
    Slide Generation Node: Creates a comprehensive PowerPoint presentation.

    This node:
    - Takes the refined plan and verified facts
    - Generates 15-20 detailed slides using GPT-4
    - Creates a downloadable PowerPoint file (.pptx)
    - Each slide supports 3-5 minutes of lecture content
    """
    start_time = time.time()

    print(f"⚙ Generating PowerPoint presentation with 15-20 slides...")

    # Prepare verified claims
    verified_claims = state.prioritized_claims

    claims_summary = ""
    for idx, claim in enumerate(verified_claims, 1):
        claims_summary += f"\n[Claim {idx}]\n"
        claims_summary += f"Content: {claim.claim}\n"
        claims_summary += f"Source: {claim.source}\n"
        claims_summary += f"Citation ID: {claim.citation_id}\n"

    # Prepare citations
    citations_list = ""
    for citation in state.citations[:20]:
        citations_list += f"\n[{citation.id}] {citation.title}\n"
        citations_list += f"    URL: {citation.url}\n"
        citations_list += f"    Snippet: {citation.snippet[:100]}...\n"

    # Load prompt template
    prompt_template = load_prompt("slides_prompt")
    prompt = prompt_template.format(
        topic=state.topic,
        refined_plan=json.dumps(state.refined_plan.dict(), indent=2) if state.refined_plan else "{}",
        verified_claims=claims_summary,
        citations=citations_list
    )

    # Call OpenAI with higher token limit for comprehensive slides
    print("  → Requesting slide content from GPT-4...")
    slides_response = call_openai(prompt, max_tokens=4000)

    try:
        # Parse JSON response
        # Remove markdown code blocks if present
        slides_response_clean = slides_response.strip()
        if slides_response_clean.startswith('```json'):
            slides_response_clean = slides_response_clean[7:]
        if slides_response_clean.startswith('```'):
            slides_response_clean = slides_response_clean[3:]
        if slides_response_clean.endswith('```'):
            slides_response_clean = slides_response_clean[:-3]

        slides_data = json.loads(slides_response_clean.strip())

        print(f"  → Parsed {len(slides_data)} slides from GPT-4 response")

        # Create output directory if it doesn't exist
        os.makedirs('outputs', exist_ok=True)

        # Generate filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_topic = "".join(c if c.isalnum() or c in (' ', '_') else '_' for c in state.topic)
        safe_topic = safe_topic.replace(' ', '_')[:50]
        filename = f"slides_{safe_topic}_{timestamp}.pptx"
        output_path = os.path.join('outputs', filename)

        # Create PowerPoint
        print(f"  → Creating PowerPoint file...")
        create_powerpoint(slides_data, state.topic, output_path)

        # Update state
        state.slides_file_path = output_path
        state.slides_data = slides_data
        state.current_node = "slide_generation"
        state.requires_human_input = False

        # Log execution
        execution_time = (time.time() - start_time) * 1000
        execution_logger.log_node_execution(
            node_name="slide_generation_node",
            inputs={
                "topic": state.topic,
                "verified_claims_count": len(verified_claims),
                "citations_count": len(state.citations)
            },
            prompt=prompt,
            output={
                "slides_count": len(slides_data),
                "file_path": output_path,
                "slides_generated": True
            },
            model_settings=get_model_settings(),
            execution_time_ms=execution_time
        )

        print(f"✓ PowerPoint generation completed!")
        print(f"  - Slides: {len(slides_data)}")
        print(f"  - File: {output_path}")

    except json.JSONDecodeError as e:
        print(f"✗ Error parsing slide data: {e}")
        print(f"Response preview: {slides_response[:200]}")
        # Store raw response for debugging
        state.slides_data = {"error": "JSON parsing failed", "raw_response": slides_response[:500]}
        state.slides_file_path = None

    return state
